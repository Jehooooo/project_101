from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color palette
DARK_BG     = RGBColor(0x0D, 0x1B, 0x2A)   # deep navy
ACCENT      = RGBColor(0x00, 0xB4, 0xD8)   # cyan
ACCENT2     = RGBColor(0x90, 0xE0, 0xEF)   # light cyan
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xCA, 0xD3, 0xE0)
CARD_BG     = RGBColor(0x1A, 0x2D, 0x42)   # slightly lighter navy
YELLOW      = RGBColor(0xFF, 0xD6, 0x00)
GREEN       = RGBColor(0x06, 0xD6, 0xA0)
RED         = RGBColor(0xFF, 0x6B, 0x6B)


def add_bg(slide, color=DARK_BG):
    from pptx.util import Inches
    bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h, font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txb


def add_multiline(slide, lines, l, t, w, h, font_size=13, color=LIGHT_GRAY,
                  bold_first=False, line_spacing=None):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"
        run.font.bold = (bold_first and i == 0)
        if line_spacing:
            from pptx.util import Pt as Pt2
            p.space_after = Pt2(line_spacing)
    return txb


def accent_bar(slide, t=0.92, h=0.07):
    bar = slide.shapes.add_shape(1, 0, Inches(t), prs.slide_width, Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()


def slide_number(slide, num):
    add_textbox(slide, f"{num:02d} / 14", 11.8, 7.1, 1.2, 0.3,
                font_size=9, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 1 – Title / Cover
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

# Left accent stripe
stripe = sl.shapes.add_shape(1, 0, 0, Inches(0.25), prs.slide_height)
stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

# Decorative right circle
circ = sl.shapes.add_shape(9, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = CARD_BG; circ.line.fill.background()

add_textbox(sl, "DMMMSU-SLUC", 0.6, 1.0, 12, 0.7, font_size=13, color=ACCENT, bold=True)
add_textbox(sl, "Disaster / Emergency Incident\nReport Monitoring System", 0.6, 1.55, 10, 2.1,
            font_size=38, bold=True, color=WHITE)
add_textbox(sl, "System Documentation & Presentation", 0.6, 3.5, 10, 0.5,
            font_size=16, color=ACCENT2, italic=True)

accent_bar(sl, t=4.3, h=0.06)

# Members block
members = [
    "Jehouse Biscarra", "Laurice Mae Ancheta", "Raymund John Prado",
    "Trisha Dacumos", "Christian Obungen", "Jezebel Subia",
    "Justin Contindenas", "Justine Jade Tavas", "John Lagao", "Cristelle Trisinio"
]
add_textbox(sl, "GROUP MEMBERS", 0.6, 4.55, 12, 0.35, font_size=10, color=ACCENT, bold=True)
col1 = members[:5]
col2 = members[5:]
for i, m in enumerate(col1):
    add_textbox(sl, f"• {m}", 0.6, 5.05 + i*0.38, 5.5, 0.38, font_size=12, color=LIGHT_GRAY)
for i, m in enumerate(col2):
    add_textbox(sl, f"• {m}", 6.5, 5.05 + i*0.38, 5.5, 0.38, font_size=12, color=LIGHT_GRAY)

slide_number(sl, 1)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 2 – Table of Contents
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "TABLE OF CONTENTS", 0.5, 0.15, 12, 0.6, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Presentation Outline", 0.5, 0.65, 12, 0.7, font_size=28, bold=True, color=WHITE)

toc_items = [
    ("01", "Introduction to the System", "About the system, its purpose and goals"),
    ("02", "Technology Stack", "Programming languages & frameworks used"),
    ("03", "Database Type & Overview", "MySQL database and its role"),
    ("04", "Database Tables", "Users & Incidents table structure"),
    ("05", "Table Relationships", "Foreign keys & relational constraints"),
    ("06", "Entity-Relationship Diagram (ERD)", "Visual representation of data model"),
    ("07", "Normalization", "1NF, 2NF, 3NF applied to the schema"),
    ("08", "Development Life Cycle Model", "SDLC methodology used"),
    ("09–14", "System Features", "Core features, roles, and workflows"),
]

for i, (num, title, desc) in enumerate(toc_items):
    row = i
    col_offset = 0 if i < 5 else 6.6
    row_idx = i if i < 5 else i - 5
    y = 1.65 + row_idx * 1.0

    card = add_rect(sl, col_offset + 0.4, y, 6.0, 0.8, CARD_BG)
    add_textbox(sl, num, col_offset + 0.55, y + 0.05, 0.7, 0.35,
                font_size=14, bold=True, color=ACCENT)
    add_textbox(sl, title, col_offset + 1.2, y + 0.05, 4.8, 0.35,
                font_size=12, bold=True, color=WHITE)
    add_textbox(sl, desc, col_offset + 1.2, y + 0.42, 4.8, 0.32,
                font_size=9, color=LIGHT_GRAY)

slide_number(sl, 2)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 3 – About the System
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "INTRODUCTION", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "About the System", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)

add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# Left column
add_rect(sl, 0.4, 1.73, 6.0, 4.8, CARD_BG)
add_textbox(sl, "📋  What is it?", 0.65, 1.88, 5.5, 0.4, font_size=13, bold=True, color=ACCENT)
desc_lines = [
    "The DMMMSU-SLUC Incident Monitoring System is a",
    "web-based application developed for Don Mariano",
    "Marcos Memorial State University – San Luis Campus.",
    "",
    "It provides a centralized, digital platform for:",
    "  • Reporting disaster & emergency incidents",
    "  • Tracking incident status in real-time",
    "  • Managing the full incident lifecycle",
    "  • Generating automated PDF reports",
    "  • Sending email & SMS notifications",
]
add_multiline(sl, desc_lines, 0.65, 2.35, 5.7, 4.0, font_size=12, color=LIGHT_GRAY)

# Right column
add_rect(sl, 6.7, 1.73, 6.2, 2.25, CARD_BG)
add_textbox(sl, "🎯  System Goal", 6.95, 1.88, 5.7, 0.4, font_size=13, bold=True, color=ACCENT)
goal_lines = [
    "Replace manual, paper-based reporting with an",
    "automated, role-based digital workflow — ensuring",
    "incidents are logged accurately, acted upon quickly,",
    "and fully documented from submission to resolution.",
]
add_multiline(sl, goal_lines, 6.95, 2.35, 5.9, 1.5, font_size=12, color=LIGHT_GRAY)

add_rect(sl, 6.7, 4.2, 6.2, 2.33, CARD_BG)
add_textbox(sl, "🔄  Incident Workflow", 6.95, 4.35, 5.7, 0.4, font_size=13, bold=True, color=ACCENT)
wf_lines = [
    "1. Staff submits incident  →  Status: Pending",
    "2. Admin reviews & acts   →  Status: In Progress",
    "3. Staff resolves issue   →  Status: Solved",
    "4. PDF auto-generated + email/SMS sent",
]
add_multiline(sl, wf_lines, 6.95, 4.82, 5.9, 1.5, font_size=12, color=LIGHT_GRAY)

slide_number(sl, 3)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 4 – Technology Stack
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "INTRODUCTION", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Technology Stack", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

tech_cols = [
    {
        "title": "⚛  Frontend",
        "color": RGBColor(0x61, 0xDB, 0xFB),
        "items": [
            "React 18",
            "TypeScript",
            "Vite (build tool)",
            "Tailwind CSS",
            "shadcn/ui components",
            "Recharts (data viz)",
            "React Router v6",
            "Axios (HTTP client)",
        ],
    },
    {
        "title": "🐍  Backend",
        "color": GREEN,
        "items": [
            "Python 3.x",
            "Flask (web framework)",
            "SQLAlchemy (ORM)",
            "Flask-JWT-Extended",
            "Flask-CORS",
            "Werkzeug (security)",
            "ReportLab (PDF gen)",
            "Twilio / SMTP (notify)",
        ],
    },
    {
        "title": "🗄  Database",
        "color": YELLOW,
        "items": [
            "MySQL 8",
            "InnoDB engine",
            "utf8mb4 charset",
            "PyMySQL connector",
            "campus_incidents_db",
            "",
            "Schema managed via",
            "SQLAlchemy models",
        ],
    },
]

for ci, col in enumerate(tech_cols):
    x = 0.4 + ci * 4.18
    card = add_rect(sl, x, 1.75, 4.0, 5.35, CARD_BG)
    # top accent strip
    top = sl.shapes.add_shape(1, Inches(x), Inches(1.75), Inches(4.0), Inches(0.18))
    top.fill.solid(); top.fill.fore_color.rgb = col["color"]; top.line.fill.background()
    add_textbox(sl, col["title"], x + 0.15, 2.0, 3.7, 0.45,
                font_size=14, bold=True, color=col["color"])
    for j, item in enumerate(col["items"]):
        add_textbox(sl, f"  {item}", x + 0.15, 2.55 + j * 0.54, 3.7, 0.48,
                    font_size=12, color=LIGHT_GRAY)

slide_number(sl, 4)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 5 – Database Type & Overview
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "DATABASE", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Database Type & Overview", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

add_rect(sl, 0.4, 1.73, 12.5, 1.7, CARD_BG)
add_textbox(sl, "🗄  MySQL Relational Database", 0.7, 1.88, 12, 0.45, font_size=15, bold=True, color=ACCENT)
add_textbox(sl,
    "The system uses MySQL 8 as its relational database management system (RDBMS). "
    "MySQL is an open-source, enterprise-grade relational database that stores data in structured tables "
    "with defined relationships, enforcing data integrity through foreign keys and constraints.",
    0.7, 2.35, 12, 0.85, font_size=12, color=LIGHT_GRAY)

props = [
    ("Database Name", "campus_incidents_db"),
    ("Engine", "InnoDB (supports FK & ACID transactions)"),
    ("Charset", "utf8mb4 — full Unicode support"),
    ("Collation", "utf8mb4_unicode_ci"),
    ("Connection", "PyMySQL via SQLAlchemy ORM"),
    ("Auth", "PBKDF2-SHA256 password hashing"),
]
for i, (k, v) in enumerate(props):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.35
    y = 3.65 + row * 1.15
    add_rect(sl, x, y, 6.0, 0.95, RGBColor(0x14, 0x24, 0x38))
    add_textbox(sl, k, x + 0.2, y + 0.07, 5.6, 0.32, font_size=10, color=ACCENT, bold=True)
    add_textbox(sl, v, x + 0.2, y + 0.43, 5.6, 0.4, font_size=12, color=WHITE)

slide_number(sl, 5)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 6 – Database Tables
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "DATABASE", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Database Tables", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# USERS table
add_rect(sl, 0.4, 1.73, 5.9, 0.42, ACCENT)
add_textbox(sl, "👤  users", 0.6, 1.78, 5.5, 0.33, font_size=13, bold=True, color=DARK_BG)
users_cols = [
    ("user_id", "INT UNSIGNED", "PK, AUTO_INCREMENT"),
    ("email", "VARCHAR(120)", "UNIQUE, NOT NULL"),
    ("password_hash", "VARCHAR(255)", "Werkzeug hash"),
    ("first_name / last_name", "VARCHAR(100)", "NOT NULL"),
    ("role", "ENUM('admin','staff')", "DEFAULT 'staff'"),
    ("phone", "VARCHAR(20)", "NULL"),
    ("is_active", "TINYINT(1)", "DEFAULT 1"),
    ("created_at / updated_at", "DATETIME", "Auto-managed timestamps"),
]
for i, (col, dtype, note) in enumerate(users_cols):
    bg_c = RGBColor(0x12, 0x22, 0x35) if i % 2 == 0 else CARD_BG
    add_rect(sl, 0.4, 2.15 + i * 0.59, 5.9, 0.57, bg_c)
    add_textbox(sl, col, 0.55, 2.2 + i * 0.59, 2.0, 0.28, font_size=10, color=ACCENT2, bold=True)
    add_textbox(sl, dtype, 2.6, 2.2 + i * 0.59, 1.7, 0.28, font_size=10, color=YELLOW)
    add_textbox(sl, note, 4.35, 2.2 + i * 0.59, 1.9, 0.28, font_size=9, color=LIGHT_GRAY)

# INCIDENTS table
add_rect(sl, 6.8, 1.73, 6.1, 0.42, GREEN)
add_textbox(sl, "📄  incidents", 7.0, 1.78, 5.8, 0.33, font_size=13, bold=True, color=DARK_BG)
inc_cols = [
    ("id", "INT UNSIGNED", "PK, AUTO_INCREMENT"),
    ("incident_id", "VARCHAR(30)", "UNIQUE identifier"),
    ("date / time", "DATE / TIME", "NOT NULL"),
    ("location", "VARCHAR(255)", "NOT NULL"),
    ("cause", "VARCHAR(255)", "NOT NULL"),
    ("description", "LONGTEXT", "NOT NULL"),
    ("file_data", "LONGBLOB", "Binary file storage"),
    ("file_name / file_mime", "VARCHAR", "For download headers"),
    ("supporting_file", "VARCHAR(255)", "Legacy FS path"),
    ("status", "ENUM(Pending…)", "NOT NULL DEFAULT Pending"),
    ("reported_by", "INT UNSIGNED", "FK → users.user_id"),
    ("created_at / updated_at", "DATETIME", "Timestamps"),
]
for i, (col, dtype, note) in enumerate(inc_cols):
    bg_c = RGBColor(0x12, 0x22, 0x35) if i % 2 == 0 else CARD_BG
    y_base = 2.15 + i * 0.44
    add_rect(sl, 6.8, y_base, 6.1, 0.42, bg_c)
    add_textbox(sl, col, 6.95, y_base + 0.05, 2.0, 0.28, font_size=9.5, color=ACCENT2, bold=True)
    add_textbox(sl, dtype, 9.0, y_base + 0.05, 1.9, 0.28, font_size=9.5, color=YELLOW)
    add_textbox(sl, note, 10.95, y_base + 0.05, 1.85, 0.28, font_size=8.5, color=LIGHT_GRAY)

slide_number(sl, 6)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 7 – Table Relationships
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "DATABASE", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Table Relationships", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# Diagram boxes
add_rect(sl, 0.5, 1.8, 4.0, 4.5, CARD_BG, line_color=ACCENT, line_width=Pt(1.5))
add_textbox(sl, "users", 0.7, 1.88, 3.6, 0.45, font_size=16, bold=True, color=ACCENT)
user_fields = ["PK  user_id", "email (UQ)", "password_hash", "first_name", "last_name",
               "role (admin/staff)", "phone", "is_active", "created_at", "updated_at"]
for i, f in enumerate(user_fields):
    add_textbox(sl, f, 0.7, 2.42 + i * 0.37, 3.6, 0.35, font_size=11,
                color=YELLOW if f.startswith("PK") else LIGHT_GRAY)

# Relationship arrow area
add_textbox(sl, "1", 4.7, 3.4, 0.5, 0.4, font_size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(sl, "─────", 4.85, 3.55, 1.2, 0.3, font_size=14, color=ACCENT)
add_textbox(sl, "∞", 6.2, 3.4, 0.5, 0.4, font_size=22, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

add_rect(sl, 0.4, 6.55, 10.0, 0.5, CARD_BG)
add_textbox(sl, "FK Constraint:  incidents.reported_by  →  users.user_id  |  ON UPDATE CASCADE  |  ON DELETE RESTRICT",
            0.6, 6.62, 9.8, 0.35, font_size=10, color=ACCENT2)

add_rect(sl, 6.8, 1.8, 5.9, 4.5, CARD_BG, line_color=GREEN, line_width=Pt(1.5))
add_textbox(sl, "incidents", 7.0, 1.88, 5.5, 0.45, font_size=16, bold=True, color=GREEN)
inc_fields = ["PK  id", "incident_id (UQ)", "date / time", "location", "cause",
              "description", "file_data (BLOB)", "status (ENUM)", "FK  reported_by", "created_at"]
for i, f in enumerate(inc_fields):
    c = YELLOW if f.startswith("PK") else (RGBColor(0xFF, 0x90, 0x40) if f.startswith("FK") else LIGHT_GRAY)
    add_textbox(sl, f, 7.0, 2.42 + i * 0.37, 5.5, 0.35, font_size=11, color=c)

rel_box = add_rect(sl, 0.4, 7.05, 12.5, 0.37, RGBColor(0x06, 0x20, 0x30))
add_textbox(sl, "Cardinality:  ONE user can report MANY incidents  (1 : N relationship)",
            0.6, 7.1, 12, 0.28, font_size=11, color=ACCENT2, bold=True)

slide_number(sl, 7)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 8 – ERD
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "DATABASE", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Entity-Relationship Diagram (ERD)", 0.5, 0.58, 12, 0.7, font_size=28, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# users entity box
ue = add_rect(sl, 0.6, 1.8, 4.2, 4.9, CARD_BG, line_color=ACCENT, line_width=Pt(2))
add_rect(sl, 0.6, 1.8, 4.2, 0.45, ACCENT)
add_textbox(sl, "USERS", 0.75, 1.85, 3.9, 0.35, font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

u_attrs = [
    ("PK", "user_id", "INT UNSIGNED"),
    ("",   "email", "VARCHAR(120)"),
    ("",   "password_hash", "VARCHAR(255)"),
    ("",   "first_name", "VARCHAR(100)"),
    ("",   "last_name", "VARCHAR(100)"),
    ("",   "role", "ENUM"),
    ("",   "phone", "VARCHAR(20)"),
    ("",   "is_active", "TINYINT"),
    ("",   "created_at", "DATETIME"),
    ("",   "updated_at", "DATETIME"),
]
for i, (tag, name, dtype) in enumerate(u_attrs):
    add_textbox(sl, tag, 0.7, 2.35 + i * 0.43, 0.4, 0.33, font_size=8, bold=True,
                color=YELLOW if tag == "PK" else LIGHT_GRAY)
    add_textbox(sl, name, 1.1, 2.35 + i * 0.43, 2.0, 0.33, font_size=10, color=WHITE)
    add_textbox(sl, dtype, 3.15, 2.35 + i * 0.43, 1.55, 0.33, font_size=9, color=LIGHT_GRAY, italic=True)

# relationship diamond
add_rect(sl, 4.95, 3.6, 1.5, 0.85, DARK_BG)
add_textbox(sl, "REPORTS", 5.0, 3.65, 1.45, 0.75, font_size=11, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
add_textbox(sl, "1", 4.68, 3.85, 0.35, 0.4, font_size=16, bold=True, color=ACCENT)
add_textbox(sl, "N", 6.5, 3.85, 0.35, 0.4, font_size=16, bold=True, color=GREEN)

# incidents entity box
add_rect(sl, 6.9, 1.8, 5.5, 4.9, CARD_BG, line_color=GREEN, line_width=Pt(2))
add_rect(sl, 6.9, 1.8, 5.5, 0.45, GREEN)
add_textbox(sl, "INCIDENTS", 7.05, 1.85, 5.15, 0.35, font_size=14, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)

inc_attrs = [
    ("PK", "id", "INT UNSIGNED"),
    ("",   "incident_id", "VARCHAR(30)"),
    ("",   "date", "DATE"),
    ("",   "time", "TIME"),
    ("",   "location", "VARCHAR(255)"),
    ("",   "cause", "VARCHAR(255)"),
    ("",   "description", "LONGTEXT"),
    ("",   "file_data", "LONGBLOB"),
    ("",   "status", "ENUM"),
    ("FK", "reported_by", "INT UNSIGNED"),
]
for i, (tag, name, dtype) in enumerate(inc_attrs):
    c = YELLOW if tag == "PK" else (RGBColor(0xFF, 0x90, 0x40) if tag == "FK" else LIGHT_GRAY)
    add_textbox(sl, tag, 6.98, 2.35 + i * 0.43, 0.4, 0.33, font_size=8, bold=True, color=c)
    add_textbox(sl, name, 7.38, 2.35 + i * 0.43, 2.5, 0.33, font_size=10, color=WHITE)
    add_textbox(sl, dtype, 9.95, 2.35 + i * 0.43, 2.35, 0.33, font_size=9, color=LIGHT_GRAY, italic=True)

slide_number(sl, 8)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 9 – Normalization
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "DATABASE", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Database Normalization", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

nf_data = [
    {
        "label": "1NF",
        "title": "First Normal Form",
        "color": RGBColor(0x00, 0xB4, 0xD8),
        "rule": "Each column must hold atomic (indivisible) values; no repeating groups.",
        "applied": [
            "Each row in users and incidents holds exactly one value per column.",
            "first_name and last_name are stored as separate columns (atomic).",
            "file_name, file_mime, file_data are individual columns — not arrays.",
            "Every row is uniquely identified by its primary key (user_id / id).",
        ],
    },
    {
        "label": "2NF",
        "title": "Second Normal Form",
        "color": GREEN,
        "rule": "Must be in 1NF; all non-key columns must fully depend on the entire primary key.",
        "applied": [
            "Both tables use a single-column primary key (PK) — no composite keys.",
            "In users: all attributes (email, role, phone…) depend solely on user_id.",
            "In incidents: all attributes depend solely on id — no partial dependency exists.",
            "incident_id is a UNIQUE business key, not part of the PK, so no partial deps.",
        ],
    },
    {
        "label": "3NF",
        "title": "Third Normal Form",
        "color": YELLOW,
        "rule": "Must be in 2NF; no non-key column depends on another non-key column (no transitive deps).",
        "applied": [
            "No derived columns exist — all data is stored, not computed from other fields.",
            "Reporter details are in users; incidents only stores the FK reported_by.",
            "status is an ENUM directly on the row — not linked via a separate status table.",
            "There are no transitive dependencies between non-key columns.",
        ],
    },
]

for i, nf in enumerate(nf_data):
    y = 1.73 + i * 1.8
    add_rect(sl, 0.4, y, 12.5, 1.72, CARD_BG)
    badge = sl.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(1.1), Inches(1.72))
    badge.fill.solid(); badge.fill.fore_color.rgb = nf["color"]; badge.line.fill.background()
    add_textbox(sl, nf["label"], 0.4, y + 0.5, 1.1, 0.65, font_size=22, bold=True,
                color=DARK_BG, align=PP_ALIGN.CENTER)
    add_textbox(sl, nf["title"], 1.65, y + 0.08, 10.8, 0.4, font_size=13, bold=True, color=nf["color"])
    add_textbox(sl, f'Rule: "{nf["rule"]}"', 1.65, y + 0.48, 10.8, 0.38, font_size=10,
                color=LIGHT_GRAY, italic=True)
    pts = "  •  " + "     •  ".join(nf["applied"])
    add_textbox(sl, "  •  " + nf["applied"][0], 1.65, y + 0.88, 10.8, 0.28, font_size=10, color=WHITE)
    add_textbox(sl, "  •  " + nf["applied"][1], 1.65, y + 1.12, 10.8, 0.28, font_size=10, color=WHITE)
    add_textbox(sl, "  •  " + nf["applied"][2], 1.65, y + 1.36, 5.2, 0.28, font_size=10, color=WHITE)
    add_textbox(sl, "  •  " + nf["applied"][3], 6.9, y + 1.36, 5.9, 0.28, font_size=10, color=WHITE)

slide_number(sl, 9)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 10 – SDLC Model
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "SYSTEM", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Development Life Cycle Model", 0.5, 0.58, 12, 0.7, font_size=28, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

add_rect(sl, 0.4, 1.73, 12.5, 0.55, CARD_BG)
add_textbox(sl, "Methodology:  Agile Software Development Model (Iterative & Incremental)",
            0.6, 1.82, 12, 0.4, font_size=13, bold=True, color=ACCENT)

phases = [
    ("Planning", ACCENT, "Define scope, goals, team roles,\nand project timeline."),
    ("Requirements\nAnalysis", GREEN, "Gather system requirements,\nuser stories, and use cases."),
    ("System\nDesign", YELLOW, "Design DB schema, ERD,\nAPI contracts, and UI mockups."),
    ("Implementation", RGBColor(0xFF, 0x90, 0x00), "Code frontend (React/TS)\nand backend (Flask/Python)."),
    ("Testing", RED, "Unit tests, integration tests,\nmanual QA, and bug fixes."),
    ("Deployment", RGBColor(0xC7, 0x7D, 0xFF), "Deploy on local server,\nreview, and release."),
]

w = 12.5 / len(phases)
for i, (phase, color, desc) in enumerate(phases):
    x = 0.4 + i * w
    add_rect(sl, x, 2.45, w - 0.06, 0.5, color)
    add_textbox(sl, phase, x + 0.05, 2.47, w - 0.12, 0.45, font_size=9.5, bold=True,
                color=DARK_BG, align=PP_ALIGN.CENTER)
    add_rect(sl, x, 3.05, w - 0.06, 3.8, CARD_BG)
    # circle step
    circ_s = sl.shapes.add_shape(9, Inches(x + (w - 0.06) / 2 - 0.3), Inches(3.25),
                                  Inches(0.6), Inches(0.6))
    circ_s.fill.solid(); circ_s.fill.fore_color.rgb = color; circ_s.line.fill.background()
    add_textbox(sl, str(i + 1), x + (w - 0.06) / 2 - 0.3, 3.25, 0.6, 0.6,
                font_size=12, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    add_textbox(sl, desc, x + 0.05, 4.0, w - 0.12, 1.8, font_size=10, color=LIGHT_GRAY,
                align=PP_ALIGN.CENTER)

add_rect(sl, 0.4, 7.0, 12.5, 0.42, CARD_BG)
add_textbox(sl,
    "Agile allows iterative delivery: each sprint produces working software, enabling rapid feedback, testing, and course-correction.",
    0.6, 7.06, 12, 0.32, font_size=10, color=ACCENT2, italic=True)

slide_number(sl, 10)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 11 – System Features (Core)
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "SYSTEM", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Core System Features", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

features = [
    ("📋", "Incident Reporting", ACCENT,
     "Staff submit incident reports with date, time, location, cause, description, and optional file attachments (DOCX/PDF)."),
    ("🔄", "Automated Status Workflow", GREEN,
     "Incidents progress through a structured pipeline: Pending → In Progress → Solved with role-based transitions."),
    ("📊", "Analytics Dashboard", YELLOW,
     "Admins view real-time data: total counts, weekly bar charts, monthly trend lines, and status pie charts."),
    ("🔍", "Search & Filtering", RGBColor(0xC7, 0x7D, 0xFF),
     "Filter incidents by date range, status, location, or keyword. Paginated results for performance."),
    ("📄", "PDF Report Generation", RED,
     "Individual incident PDFs and compiled full reports are auto-generated at each workflow step using ReportLab."),
    ("🔔", "Notifications", RGBColor(0xFF, 0x90, 0x00),
     "Automated email (SMTP) and SMS (Twilio) notifications sent on new submissions and status updates."),
]

for i, (icon, title, color, desc) in enumerate(features):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.35
    y = 1.73 + row * 1.88
    add_rect(sl, x, y, 6.1, 1.75, CARD_BG)
    top = sl.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.18), Inches(1.75))
    top.fill.solid(); top.fill.fore_color.rgb = color; top.line.fill.background()
    add_textbox(sl, icon, x + 0.28, y + 0.12, 0.55, 0.55, font_size=22, color=color)
    add_textbox(sl, title, x + 0.88, y + 0.14, 5.1, 0.45, font_size=13, bold=True, color=WHITE)
    add_textbox(sl, desc, x + 0.28, y + 0.68, 5.7, 0.97, font_size=10.5, color=LIGHT_GRAY)

slide_number(sl, 11)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 12 – User Roles & Access Control
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "SYSTEM", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "User Roles & Access Control", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# Staff card
add_rect(sl, 0.4, 1.73, 6.0, 5.35, CARD_BG, line_color=ACCENT, line_width=Pt(1.5))
add_rect(sl, 0.4, 1.73, 6.0, 0.5, ACCENT)
add_textbox(sl, "👤  Staff Role", 0.6, 1.8, 5.6, 0.38, font_size=15, bold=True, color=DARK_BG)
staff_perms = [
    ("✅", "Log in to the staff portal"),
    ("✅", "Submit new incident reports"),
    ("✅", "Upload supporting documents"),
    ("✅", "View only their own reports"),
    ("✅", "Mark incidents as Solved"),
    ("✅", "Download PDF of own reports"),
    ("❌", "View other staff's incidents"),
    ("❌", "Access analytics dashboard"),
    ("❌", "Manage user accounts"),
]
for i, (icon, perm) in enumerate(staff_perms):
    add_textbox(sl, f"{icon}  {perm}", 0.6, 2.35 + i * 0.5, 5.7, 0.45, font_size=11.5, color=WHITE)

# Admin card
add_rect(sl, 6.8, 1.73, 6.1, 5.35, CARD_BG, line_color=GREEN, line_width=Pt(1.5))
add_rect(sl, 6.8, 1.73, 6.1, 0.5, GREEN)
add_textbox(sl, "🛡  Administrator Role", 7.0, 1.8, 5.8, 0.38, font_size=15, bold=True, color=DARK_BG)
admin_perms = [
    ("✅", "Log in to the admin portal"),
    ("✅", "View ALL incident reports"),
    ("✅", "Update status to In Progress"),
    ("✅", "Access analytics dashboard"),
    ("✅", "Generate compiled PDF reports"),
    ("✅", "Create / edit / deactivate users"),
    ("✅", "Manage notification settings"),
    ("✅", "Delete any incident"),
    ("✅", "Full system access"),
]
for i, (icon, perm) in enumerate(admin_perms):
    add_textbox(sl, f"{icon}  {perm}", 7.0, 2.35 + i * 0.5, 5.9, 0.45, font_size=11.5, color=WHITE)

slide_number(sl, 12)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 13 – Project Structure
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)
accent_bar(sl)

add_textbox(sl, "SYSTEM", 0.5, 0.15, 12, 0.4, font_size=11, color=ACCENT, bold=True)
add_textbox(sl, "Project Structure", 0.5, 0.58, 12, 0.7, font_size=30, bold=True, color=WHITE)
add_rect(sl, 0.4, 1.52, 12.5, 0.06, ACCENT)

# Backend tree
add_rect(sl, 0.4, 1.73, 5.9, 5.35, CARD_BG)
add_rect(sl, 0.4, 1.73, 5.9, 0.4, ACCENT)
add_textbox(sl, "📁  backend/", 0.6, 1.78, 5.5, 0.3, font_size=12, bold=True, color=DARK_BG)
be_tree = [
    ("", "app.py", "Main Flask API — all routes & logic"),
    ("", "schema.sql", "MySQL DDL schema definitions"),
    ("", "mysql_setup.py", "DB init & seeding script"),
    ("", ".env.example", "Environment variables template"),
    ("", "requirements.txt", "Python dependencies list"),
    ("├── models/", "", ""),
    ("│   ├──", "database.py", "SQLAlchemy db instance & init"),
    ("│   ├──", "user.py", "User model & to_dict()"),
    ("│   └──", "incident.py", "Incident model & helpers"),
    ("└── utils/", "", ""),
    ("    ├──", "pdf_generator.py", "ReportLab PDF generation"),
    ("    └──", "notifications.py", "Email & SMS dispatch"),
]
for i, (prefix, name, desc) in enumerate(be_tree):
    if not name:
        add_textbox(sl, prefix, 0.6, 2.22 + i * 0.38, 5.6, 0.35, font_size=10.5,
                    color=ACCENT2, bold=True)
    else:
        add_textbox(sl, f"{prefix}", 0.6, 2.22 + i * 0.38, 1.1, 0.35, font_size=10, color=LIGHT_GRAY)
        add_textbox(sl, name, 1.7, 2.22 + i * 0.38, 1.8, 0.35, font_size=10, bold=True, color=YELLOW)
        add_textbox(sl, f"— {desc}", 3.55, 2.22 + i * 0.38, 2.65, 0.35, font_size=9, color=LIGHT_GRAY)

# Frontend tree
add_rect(sl, 6.75, 1.73, 6.15, 5.35, CARD_BG)
add_rect(sl, 6.75, 1.73, 6.15, 0.4, GREEN)
add_textbox(sl, "📁  app/ (frontend)", 6.95, 1.78, 5.8, 0.3, font_size=12, bold=True, color=DARK_BG)
fe_tree = [
    ("", "package.json", "Dependencies & scripts"),
    ("", "vite.config.ts", "Vite build configuration"),
    ("", "index.html", "HTML entry point"),
    ("└── src/", "", ""),
    ("    ├──", "App.tsx", "Root component & routing"),
    ("    ├──", "main.tsx", "React DOM entry"),
    ("    ├── pages/", "", ""),
    ("    │  ├──", "AdminDashboard.tsx", "Analytics & admin view"),
    ("    │  ├──", "IncidentReports.tsx", "Incidents list & CRUD"),
    ("    │  ├──", "StaffDashboard.tsx", "Staff home view"),
    ("    │  ├──", "UserManagement.tsx", "Admin user CRUD"),
    ("    │  ├──", "Settings.tsx", "System settings page"),
    ("    │  └──", "AdminLogin / StaffLogin", "Separate login portals"),
    ("    ├── components/", "", ""),
    ("    ├── contexts/ (AuthContext)", "", ""),
    ("    └── services/ (api.ts)", "", ""),
]
for i, (prefix, name, desc) in enumerate(fe_tree):
    if not name:
        add_textbox(sl, prefix, 6.95, 2.22 + i * 0.32, 5.8, 0.3, font_size=9.5,
                    color=ACCENT2, bold=True)
    else:
        add_textbox(sl, f"{prefix}", 6.95, 2.22 + i * 0.32, 1.4, 0.3, font_size=9.5, color=LIGHT_GRAY)
        add_textbox(sl, name, 8.4, 2.22 + i * 0.32, 2.2, 0.3, font_size=9.5, bold=True, color=YELLOW)
        if desc:
            add_textbox(sl, f"— {desc}", 10.65, 2.22 + i * 0.32, 2.2, 0.3, font_size=8.5, color=LIGHT_GRAY)

slide_number(sl, 13)

# ──────────────────────────────────────────────────────────────────────
# SLIDE 14 – Summary / Thank You
# ──────────────────────────────────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl)

stripe = sl.shapes.add_shape(1, 0, 0, Inches(0.25), prs.slide_height)
stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

circ = sl.shapes.add_shape(9, Inches(9.5), Inches(4), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = CARD_BG; circ.line.fill.background()

accent_bar(sl, t=2.7, h=0.07)

add_textbox(sl, "DMMMSU-SLUC", 0.6, 0.8, 12, 0.55, font_size=13, color=ACCENT, bold=True)
add_textbox(sl, "Summary", 0.6, 1.3, 12, 0.7, font_size=40, bold=True, color=WHITE)

summary_pts = [
    ("🏫", "System", "Web-based Disaster/Emergency Incident Report Monitoring System for DMMMSU-SLUC"),
    ("💻", "Frontend", "React 18 + TypeScript + Vite + Tailwind CSS + shadcn/ui"),
    ("🐍", "Backend", "Python Flask + SQLAlchemy + JWT Authentication + ReportLab"),
    ("🗄", "Database", "MySQL 8 (campus_incidents_db) — 2 tables, fully normalized to 3NF"),
    ("🔄", "SDLC", "Agile iterative model with 6 phases"),
    ("🔑", "Roles", "Staff (report & resolve) | Admin (manage, analyze, oversee)"),
]
for i, (icon, label, text) in enumerate(summary_pts):
    y = 3.05 + i * 0.67
    add_rect(sl, 0.5, y, 12.3, 0.6, CARD_BG)
    add_textbox(sl, icon, 0.65, y + 0.08, 0.45, 0.44, font_size=16, color=ACCENT)
    add_textbox(sl, label, 1.18, y + 0.1, 1.3, 0.38, font_size=11, bold=True, color=ACCENT)
    add_textbox(sl, text, 2.55, y + 0.1, 10.15, 0.38, font_size=11, color=WHITE)

add_textbox(sl, "Thank You!", 0.6, 7.0, 12.3, 0.42, font_size=22, bold=True,
            color=ACCENT, align=PP_ALIGN.CENTER)

slide_number(sl, 14)

# ──────────────────────────────────────────────────────────────────────
out_path = r"c:\Users\Jeho\Downloads\project_102\DMMMSU_SLUC_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
